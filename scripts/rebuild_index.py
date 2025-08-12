#!/usr/bin/env python3
"""
Rebuild a FAISS index from GCS documents and upload artifacts back to GCS.

Usage:
    python3 scripts/rebuild_index.py wyndham \
        --bucket civreply-data \
        --docs-prefix policies/{slug} \
        --index-prefix policies/{slug}/_hash_index \
        --chunk-size 1200 --overlap 180

Env vars required:
    OPENAI_API_KEY
    GOOGLE_APPLICATIONS_CREDENTIALS (path to GCP service account JSON)

Notes:
- Supports: .pdf, .docx, .pptx, .xlsx, .xls, .csv, .rtf, .txt
- Writes index.faiss / index.pkl (+ index_meta.json) to the index prefix
"""
import os, io, csv, json, tempfile, datetime, argparse, sys
from pathlib import Path
from typing import Iterable, List

from google.cloud import storage
from PyPDF2 import PdfReader
import pandas as pd
from pptx import Presentation
import docx
from striprtf.striprtf import rtf_to_text

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS

SUPPORTED_EXTS = (".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".rtf", ".txt")

# ---------- extractors ----------

def _join(lines: Iterable[str]) -> str:
    return "\n".join([str(x).strip() for x in lines if str(x).strip()])

def extract_pdf(b: bytes) -> str:
    try:
        return _join((p.extract_text() or "") for p in PdfReader(io.BytesIO(b)).pages)
    except Exception:
        return ""

def extract_docx_(b: bytes) -> str:
    try:
        d = docx.Document(io.BytesIO(b))
        return _join(p.text for p in d.paragraphs)
    except Exception:
        return ""

def extract_pptx(b: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(b))
        out = []
        for s in prs.slides:
            chunk = []
            for shp in s.shapes:
                if hasattr(shp, "text"):
                    chunk.append(shp.text)
            out.append(_join(chunk))
        return _join(out)
    except Exception:
        return ""

def extract_xlsx(b: bytes) -> str:
    try:
        frames = []
        with pd.ExcelFile(io.BytesIO(b)) as xls:
            for sheet in xls.sheet_names:
                df = xls.parse(sheet)
                frames.append(f"[Sheet: {sheet}]\n" + df.to_string(index=False))
        return "\n\n".join(frames)
    except Exception:
        return ""

def extract_xls(b: bytes) -> str:
    try:
        df = pd.read_excel(io.BytesIO(b), engine="xlrd", sheet_name=None)
        return "\n\n".join(f"[Sheet: {s}]\n{d.to_string(index=False)}" for s, d in df.items())
    except Exception:
        return ""

def extract_csv(b: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            txt = b.decode(enc, errors="ignore")
            break
        except Exception:
            pass
    return _join(", ".join(r) for r in csv.reader(io.StringIO(txt)))

def extract_rtf(b: bytes) -> str:
    try:
        return rtf_to_text(b.decode("utf-8", errors="ignore"))
    except Exception:
        return ""

def extract_txt(b: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return b.decode(enc, errors="ignore")
        except Exception:
            pass
    return ""

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx_,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls": extract_xls,
    ".csv": extract_csv,
    ".rtf": extract_rtf,
    ".txt": extract_txt,
}

# ---------- cli ----------

parser = argparse.ArgumentParser(description="Rebuild FAISS from GCS docs and upload to GCS")
parser.add_argument("slug", help="Council slug, e.g., wyndham")
parser.add_argument("--bucket", default=os.getenv("GCS_BUCKET", "civreply-data"))
parser.add_argument("--docs-prefix", default=os.getenv("DOC_PREFIX", "policies/{slug}"),
                    help="GCS prefix where documents live; {slug} will be replaced")
parser.add_argument("--index-prefix", default=os.getenv("INDEX_PREFIX", "policies/{slug}/_hash_index"),
                    help="GCS prefix where FAISS artifacts are stored; {slug} will be replaced")
parser.add_argument("--chunk-size", type=int, default=int(os.getenv("CHUNK_SIZE", 1200)))
parser.add_argument("--overlap", type=int, default=int(os.getenv("CHUNK_OVERLAP", 180)))
parser.add_argument("--embedding-model", default=os.getenv("EMBEDDING_MODEL", "text-embedding-3-small"),
                    help="OpenAI embedding model to use (e.g., text-embedding-3-small)")
parser.add_argument("--embed-batch", type=int, default=int(os.getenv("EMBED_BATCH", "64")),
                    help="Max texts per embeddings API call")
args = parser.parse_args()

if not os.environ.get("OPENAI_API_KEY"):
    sys.exit("Missing OPENAI_API_KEY")

# Resolve prefixes
args.docs_prefix = args.docs_prefix.format(slug=args.slug).rstrip("/")
args.index_prefix = args.index_prefix.format(slug=args.slug).rstrip("/")

print(f"→ Using bucket: {args.bucket}")
print(f"→ Docs prefix: gs://{args.bucket}/{args.docs_prefix}")
print(f"→ Index out : gs://{args.bucket}/{args.index_prefix}")

cli = storage.Client()

# ---------- list blobs ----------

def list_supported_blobs(cli: storage.Client, bucket: str, prefix: str):
    prefix = prefix.rstrip("/") + "/"
    return [b for b in cli.list_blobs(bucket, prefix=prefix)
            if b.name.lower().endswith(SUPPORTED_EXTS)]

blobs = list_supported_blobs(cli, args.bucket, args.docs_prefix)
if not blobs:
    sys.exit(f"No docs found under gs://{args.bucket}/{args.docs_prefix}")
print(f"→ Found {len(blobs)} file(s)")

# ---------- chunk builder ----------

def build_chunks_from_blobs(blobs) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=args.chunk_size,
                                              chunk_overlap=args.overlap)
    chunks: List[str] = []
    for i, b in enumerate(blobs, 1):
        try:
            data = b.download_as_bytes()
            ext = Path(b.name.lower()).suffix
            text = EXTRACTORS.get(ext, lambda _: "")(data)
            if not text.strip():
                print(f"   [skip no-text] {b.name}")
                continue
            header = f"[SOURCE] gs://{args.bucket}/{b.name}\n"
            chunks.extend(splitter.split_text(header + text))
            if i % 20 == 0:
                print(f"   …parsed {i}")
        except Exception as e:
            print(f"   [skip error] {b.name}: {e}")
    return chunks

print("→ Extracting & chunking…")
chunks = build_chunks_from_blobs(blobs)
if not chunks:
    sys.exit("No text chunks created — nothing to index.")
print(f"→ Built {len(chunks)} chunks")

# ---------- build embeddings + FAISS ----------

print("→ Embedding & building FAISS…")
import time
emb = OpenAIEmbeddings(
    model=args.embedding_model,
    openai_api_key=os.environ["OPENAI_API_KEY"],
)
vs = None
batch = max(1, args.embed_batch)
total = len(chunks)
for i in range(0, total, batch):
    sub = chunks[i:i+batch]
    for attempt in range(1, 7):  # up to 6 attempts with exponential backoff
        try:
            if vs is None:
                vs = FAISS.from_texts(sub, emb)
            else:
                vs.add_texts(sub)
            break
        except Exception as e:
            if attempt == 6:
                raise
            wait = min(60, 2 ** attempt)
            print(f"   [retry {attempt}/6] embedding {i+1}-{min(i+len(sub), total)} failed: {e}. Sleeping {wait}s…")
            time.sleep(wait)
    if ((i // batch) % 10) == 0 or i + batch >= total:
        print(f"   …embedded {min(i+batch, total)}/{total}")
if vs is None:
    sys.exit("No vectors computed.")

# ---------- save & upload ----------

print("→ Saving artifacts locally…")
tmp = Path(tempfile.mkdtemp()) / "faiss"
tmp.mkdir(parents=True, exist_ok=True)
vs.save_local(str(tmp))

print("→ Uploading to GCS…")
from google.api_core.retry import Retry
retry = Retry(predicate=Retry.if_transient_error, initial=1.0, maximum=60.0, multiplier=2.0, deadline=900.0)
bucket = cli.bucket(args.bucket)

# Clear previous index
for b in list(bucket.list_blobs(prefix=f"{args.index_prefix}/")):
    b.delete()

for p in tmp.rglob("*"):
    if p.is_file():
        key = f"{args.index_prefix}/{p.relative_to(tmp).as_posix()}"
        bl = bucket.blob(key)
        bl.chunk_size = 8 * 1024 * 1024  # 8MB
        bl.upload_from_filename(str(p), timeout=300, retry=retry)

meta = {
    "generated_at_utc": datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S"),
    "council": args.slug,
    "chunks": len(chunks),
    "docs_prefix": args.docs_prefix,
    "index_prefix": args.index_prefix,
    "embedding_model": args.embedding_model,
}
bucket.blob(f"{args.index_prefix}/index_meta.json").upload_from_string(
    json.dumps(meta, indent=2).encode("utf-8"), content_type="application/json"
)

print("✓ Done. Uploaded index.faiss / index.pkl + index_meta.json")
