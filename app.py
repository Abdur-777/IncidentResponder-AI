# =========================
# IncidentResponder AI — One app for Basic, Premium, Enterprise
# - Basic: KB Q&A over council docs (FAISS), branding, logs
# - Premium: Inbox triage + policy-grounded draft + approve & send
# - Enterprise: extras (signed URL TTL, data banner, larger uploads, mini analytics)
# Admin login REQUIRED to access Premium/Enterprise UI even if enabled via env.
# =========================

import os, io, re, csv, json, shutil, tempfile, datetime
from pathlib import Path

# ---- Load GCP creds from env JSON (Render/Cloud) ----
if "GCP_SA_JSON" in os.environ:
    with open("/tmp/gcs-key.json", "w") as f:
        f.write(os.environ["GCP_SA_JSON"])
    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = "/tmp/gcs-key.json"

import streamlit as st
from dotenv import load_dotenv
from google.cloud import storage

# File parsing
from PyPDF2 import PdfReader
import docx  # python-docx
from pptx import Presentation
import pandas as pd
from striprtf.striprtf import rtf_to_text

# LangChain / OpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import FAISS

# Premium (email)
import imaplib, email
from email.header import decode_header, make_header
import yagmail
from langdetect import detect as lang_detect

# =========================
#  ENV / CONFIG
# =========================
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    st.error("❌ Missing OPENAI_API_KEY")
    st.stop()

# Core (single-council deployment)
COUNCIL        = os.getenv("COUNCIL", "wyndham")
GCS_BUCKET     = os.getenv("GCS_BUCKET", "civreply-data")
DOC_PREFIX     = os.getenv("DOC_PREFIX", f"policies/{COUNCIL}")
INDEX_PREFIX   = os.getenv("INDEX_PREFIX", f"indexes/{COUNCIL}/faiss")
INDEX_ON_START = os.getenv("INDEX_ON_START", "false").lower() == "true"

# Branding / UX
PUBLIC_MODE    = os.getenv("PUBLIC_MODE", "false").lower() == "true"    # read-only portal
ADMIN_PASSWORD = os.getenv("ADMIN_PASSWORD", "")
BRAND_LOGO_URL = os.getenv("BRAND_LOGO_URL", "")
BRAND_PRIMARY  = os.getenv("BRAND_PRIMARY", "#0055a5")

# Plans (env toggles) — still require admin login to be usable
BASIC_MODE       = os.getenv("BASIC_MODE", "true").lower() == "true"
PREMIUM_MODE     = os.getenv("PREMIUM_MODE", "false").lower() == "true"  # inbox triage + draft
ENTERPRISE_MODE  = os.getenv("ENTERPRISE_MODE", "false").lower() == "true"

# Enterprise extras
SIGNED_URL_TTL_MINUTES = int(os.getenv("SIGNED_URL_TTL_MINUTES", "120"))
DATA_REGION            = os.getenv("DATA_REGION", "")       # e.g., australia-southeast1
MAX_UPLOAD_MB          = int(os.getenv("MAX_UPLOAD_MB", "50"))

# Premium email
IMAP_HOST   = os.getenv("IMAP_HOST", "")
IMAP_USER   = os.getenv("IMAP_USER", "")
IMAP_PASS   = os.getenv("IMAP_PASS", "")
IMAP_FOLDER = os.getenv("IMAP_FOLDER", "INBOX")
REPLY_FROM  = os.getenv("REPLY_FROM", "")
YAGMAIL_PASS= os.getenv("YAGMAIL_PASS", "")
EMAIL_FROM_NAME = os.getenv("EMAIL_FROM_NAME", f"{COUNCIL.title()} AI Assistant")

# Optional routing
ROUTES = {
    "waste":    os.getenv("ROUTE_WASTE",   ""),
    "roads":    os.getenv("ROUTE_ROADS",   ""),
    "animals":  os.getenv("ROUTE_ANIMALS", ""),
    "noise":    os.getenv("ROUTE_NOISE",   ""),
    "rates":    os.getenv("ROUTE_RATES",   ""),
    "parks":    os.getenv("ROUTE_PARKS",   ""),
    "planning": os.getenv("ROUTE_PLANNING",""),
    "health":   os.getenv("ROUTE_HEALTH",  ""),
}

# Paths
LOCAL_INDEX = Path(f"index/{COUNCIL}")

SUPPORTED_EXTS = (".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".rtf", ".txt")

# =========================
#  PAGE SETUP
# =========================
st.set_page_config(page_title="IncidentResponder AI", page_icon="🚨", layout="wide")

def gcs_client():
    return storage.Client()

def latest_index_time_str():
    try:
        blobs = list(gcs_client().list_blobs(GCS_BUCKET, prefix=f"{INDEX_PREFIX}/"))
        if not blobs:
            return "N/A"
        newest = max((b.updated for b in blobs if b.updated), default=None)
        if newest:
            return newest.strftime("%Y-%m-%d %H:%M UTC")
    except Exception:
        pass
    return "N/A"

# Header
with st.container():
    col1, col2 = st.columns([6, 4])
    with col1:
        if BRAND_LOGO_URL:
            st.markdown(
                f"""
                <div style="display:flex;align-items:center;gap:12px;margin-top:6px;">
                    <img src="{BRAND_LOGO_URL}" style="height:42px"/>
                    <div style="font-size:26px;font-weight:800;color:{BRAND_PRIMARY};">
                        {COUNCIL.title()} — IncidentResponder AI
                    </div>
                </div>
                """, unsafe_allow_html=True
            )
        else:
            st.markdown(
                f"<div style='font-size:26px;font-weight:800;color:{BRAND_PRIMARY};margin:6px 0'>"
                f"{COUNCIL.title()} — IncidentResponder AI</div>",
                unsafe_allow_html=True
            )
        sub = "Public read-only mode enabled." if PUBLIC_MODE else "Staff assistant with council policy knowledge."
        st.caption(sub + ("  •  " + f"Data region: {DATA_REGION}" if ENTERPRISE_MODE and DATA_REGION else ""))
    with col2:
        st.markdown(
            f"""
            <div style="text-align:right;margin-top:4px;">
                <div style="font-size:13px;color:#666;">Last indexed</div>
                <div style="font-size:15px;font-weight:600;">{latest_index_time_str()}</div>
            </div>
            """, unsafe_allow_html=True
        )

st.divider()

# =========================
#  HELPERS: GCS / INDEX
# =========================
def list_supported_blobs():
    blobs = list(gcs_client().list_blobs(GCS_BUCKET, prefix=DOC_PREFIX.rstrip("/") + "/"))
    files = [b for b in blobs if b.name.lower().endswith(SUPPORTED_EXTS)]
    files.sort(key=lambda b: b.updated or 0, reverse=True)
    return files

def download_dir_from_gcs(prefix: str, dest_dir: Path) -> bool:
    blobs = list(gcs_client().list_blobs(GCS_BUCKET, prefix=prefix.rstrip("/") + "/"))
    if not blobs:
        return False
    tmp = Path(tempfile.mkdtemp())
    for b in blobs:
        rel = b.name[len(prefix)+1:]
        if not rel:
            continue
        p = tmp / rel
        p.parent.mkdir(parents=True, exist_ok=True)
        b.download_to_filename(str(p))
    dest_dir.mkdir(parents=True, exist_ok=True)
    for p in tmp.rglob("*"):
        if p.is_file():
            tgt = dest_dir / p.relative_to(tmp)
            tgt.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(p), str(tgt))
    return True

def upload_dir_to_gcs(src_dir: Path, prefix: str):
    bucket = gcs_client().bucket(GCS_BUCKET)
    for p in src_dir.rglob("*"):
        if p.is_file():
            rel = p.relative_to(src_dir).as_posix()
            key = f"{prefix.rstrip('/')}/{rel}"
            bucket.blob(key).upload_from_filename(str(p))

def signed_url_for_blob_path(gcs_path: str, minutes: int) -> str:
    # gcs_path can be "gs://bucket/key" or just "key"
    if gcs_path.startswith("gs://"):
        _, rest = gcs_path.split("gs://", 1)
        bucket_name, key = rest.split("/", 1)
    else:
        bucket_name, key = GCS_BUCKET, gcs_path
    try:
        blob = gcs_client().bucket(bucket_name).blob(key)
        return blob.generate_signed_url(
            expiration=datetime.timedelta(minutes=minutes),
            method="GET"
        )
    except Exception:
        return ""

# =========================
#  TEXT EXTRACTORS
# =========================
def _join(lines): return "\n".join([str(x).strip() for x in lines if str(x).strip()])

def extract_pdf(data: bytes) -> str:
    try:
        pages = PdfReader(io.BytesIO(data)).pages
        return _join(p.extract_text() or "" for p in pages)
    except Exception: return ""

def extract_docx(data: bytes) -> str:
    try:
        d = docx.Document(io.BytesIO(data))
        return _join(p.text for p in d.paragraphs)
    except Exception: return ""

def extract_pptx(data: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(data))
        out = []
        for s in prs.slides:
            txts = []
            for shp in s.shapes:
                if hasattr(shp, "text"):
                    txts.append(shp.text)
            out.append(_join(txts))
        return _join(out)
    except Exception: return ""

def extract_xlsx(data: bytes) -> str:
    try:
        with pd.ExcelFile(io.BytesIO(data)) as xls:
            frames = []
            for sheet in xls.sheet_names:
                df = xls.parse(sheet)
                frames.append(f"[Sheet: {sheet}]\n" + df.to_string(index=False))
            return "\n\n".join(frames)
    except Exception: return ""

def extract_xls(data: bytes) -> str:
    try:
        df = pd.read_excel(io.BytesIO(data), engine="xlrd", sheet_name=None)
        return "\n\n".join(f"[Sheet: {s}]\n{d.to_string(index=False)}" for s, d in df.items())
    except Exception: return ""

def extract_csv(data: bytes) -> str:
    try:
        for enc in ("utf-8","utf-8-sig","latin-1"):
            try:
                txt = data.decode(enc, errors="ignore"); break
            except Exception: pass
        return _join(", ".join(r) for r in csv.reader(io.StringIO(txt)))
    except Exception: return ""

def extract_rtf(data: bytes) -> str:
    try:
        return rtf_to_text(data.decode("utf-8", errors="ignore"))
    except Exception: return ""

def extract_txt(data: bytes) -> str:
    for enc in ("utf-8","utf-8-sig","latin-1"):
        try: return data.decode(enc, errors="ignore")
        except Exception: pass
    return ""

EXTRACTORS = {
    ".pdf": extract_pdf, ".docx": extract_docx, ".pptx": extract_pptx,
    ".xlsx": extract_xlsx, ".xls": extract_xls, ".csv": extract_csv,
    ".rtf": extract_rtf, ".txt": extract_txt,
}
def extract_any(name: str, data: bytes) -> str:
    fn = EXTRACTORS.get(Path(name.lower()).suffix)
    return fn(data) if fn else ""

# =========================
#  INDEX BUILD / LOAD
# =========================
def build_index_from_gcs() -> bool:
    st.session_state["build_logs"] = []
    blobs = list_supported_blobs()
    if not blobs:
        st.session_state["build_logs"].append(f"No files under gs://{GCS_BUCKET}/{DOC_PREFIX}")
        return False

    splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=180)
    texts = []
    st.session_state["build_logs"].append(f"Indexing {len(blobs)} files…")
    for i, b in enumerate(blobs, 1):
        try:
            data = b.download_as_bytes()
        except Exception as e:
            st.session_state["build_logs"].append(f"[skip] {b.name} download failed: {e}")
            continue
        body = extract_any(b.name, data)
        if not body.strip():
            st.session_state["build_logs"].append(f"[skip] {b.name} (no text)")
            continue
        header = f"[SOURCE] gs://{GCS_BUCKET}/{b.name}\n"
        for chunk in splitter.split_text(header + body):
            texts.append(chunk)
        if i % 20 == 0:
            st.session_state["build_logs"].append(f"... parsed {i} files")

    if not texts:
        st.session_state["build_logs"].append("No text chunks produced.")
        return False

    embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
    vs = FAISS.from_texts(texts, embeddings)
    LOCAL_INDEX.parent.mkdir(parents=True, exist_ok=True)
    vs.save_local(str(LOCAL_INDEX))
    upload_dir_to_gcs(LOCAL_INDEX, INDEX_PREFIX)

    meta = {
        "council": COUNCIL,
        "chunks": len(texts),
        "generated_at_utc": datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")
    }
    try:
        storage.Client().bucket(GCS_BUCKET).blob(f"{INDEX_PREFIX}/index_meta.json")\
            .upload_from_string(json.dumps(meta, indent=2).encode("utf-8"),
                                content_type="application/json")
    except Exception:
        pass

    st.session_state["build_logs"].append(f"[done] Index uploaded to gs://{GCS_BUCKET}/{INDEX_PREFIX}")
    return True

def load_vectorstore():
    if LOCAL_INDEX.exists():
        try:
            emb = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
            return FAISS.load_local(str(LOCAL_INDEX), emb, allow_dangerous_deserialization=True)
        except Exception as e:
            st.warning(f"Local FAISS present but failed to load: {e}")

    if download_dir_from_gcs(INDEX_PREFIX, LOCAL_INDEX):
        try:
            emb = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
            return FAISS.load_local(str(LOCAL_INDEX), emb, allow_dangerous_deserialization=True)
        except Exception as e:
            st.warning(f"GCS FAISS downloaded but failed to load: {e}")

    if INDEX_ON_START:
        with st.spinner("Building index from GCS… (first boot may take minutes)"):
            ok = build_index_from_gcs()
        if ok:
            emb = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
            return FAISS.load_local(str(LOCAL_INDEX), emb, allow_dangerous_deserialization=True)

    return None

def get_llm():
    return ChatOpenAI(model_name="gpt-4o-mini", temperature=0, openai_api_key=OPENAI_API_KEY)

# =========================
#  LOGGING
# =========================
def log_query_to_gcs(mode: str, question: str, answer: str, sources: list):
    try:
        ts = datetime.datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")
        day = datetime.datetime.utcnow().strftime("%Y%m%d")
        key = f"logs/{COUNCIL}/queries-{day}.csv"
        bucket = gcs_client().bucket(GCS_BUCKET)
        blob = bucket.blob(key)

        src_str = " | ".join(sources)[:2000] if sources else ""
        def q(x): return '"' + str(x).replace('"', "'") + '"'
        line = f"{q(ts)},{q(mode)},{q(question)},{q(answer[:2000])},{q(src_str)}\n"

        if blob.exists():
            existing = blob.download_as_text()
            updated = existing + line
        else:
            header = "timestamp_utc,mode,question,answer_preview,sources\n"
            updated = header + line

        blob.upload_from_string(updated, content_type="text/csv")
    except Exception:
        pass

# =========================
#  CATEGORY DETECTION (routing)
# =========================
CATEGORY_RULES = [
    ("waste|bin|hard[-_ ]?rubbish|recycling|garbage|litter|landfill|compost|green waste|organic", "waste"),
    ("road|traffic|pothole|parking|transport|footpath|foot-path|bike|bicycle|bridge|intersection|roundabout", "roads"),
    ("animal|pet|dog|cat|livestock|impound|regist|wildlife|vet|microchip", "animals"),
    ("noise|nuisance|amenity|disturbance|music volume|loud", "noise"),
    ("rate|payment|valuation|fine|infringement|penalty|fees|charges|levy", "rates"),
    ("park|tree|environment|reserve|open[-_ ]space|green|garden|bushland|playground", "parks"),
    ("plan|permit|building|construction|overlay|zoning|town[-_ ]plan|development application|\\bDA\\b", "planning"),
    ("health|food|safety|public health|inspection|disease|mosquito|covid|hygiene", "health"),
]
def detect_category(text: str) -> str:
    t = (text or "").lower()
    for pat, cat in CATEGORY_RULES:
        if re.search(pat, t):
            return cat
    return "uncategorized"

# =========================
#  UI: SIDEBAR (Admin, Premium, Enterprise)
# =========================
st.sidebar.title("Menu")

is_admin = st.session_state.get("is_admin", False)
if not PUBLIC_MODE:
    if not is_admin:
        pwd = st.sidebar.text_input("Admin password", type="password")
        if st.sidebar.button("Login"):
            if ADMIN_PASSWORD and pwd == ADMIN_PASSWORD:
                st.session_state["is_admin"] = True
                is_admin = True
                st.sidebar.success("Admin unlocked.")
            else:
                st.sidebar.error("Incorrect password.")

    # Admin utilities (always visible to admin)
    if is_admin:
        st.sidebar.markdown("### Admin")
        if st.sidebar.button("Rebuild KB Index Now"):
            with st.spinner("Rebuilding index…"):
                ok = build_index_from_gcs()
            st.sidebar.success("Done." if ok else "No files or failed.")
        day = datetime.datetime.utcnow().strftime("%Y%m%d")
        log_key = f"logs/{COUNCIL}/queries-{day}.csv"
        if gcs_client().bucket(GCS_BUCKET).blob(log_key).exists():
            href = signed_url_for_blob_path(f"gs://{GCS_BUCKET}/{log_key}", minutes=60)
            if href: st.sidebar.markdown(f"[Download today’s log]({href})")

        # Premium features are only shown if enabled AND admin
        if PREMIUM_MODE:
            st.sidebar.markdown("### Inbox (Premium)")
            def _decode_header(raw):
                try: return str(make_header(decode_header(raw or "")))
                except Exception: return raw or ""

            def fetch_recent_emails(limit=25, since_days=7):
                out = []
                if not (IMAP_HOST and IMAP_USER and IMAP_PASS):
                    return out
                try:
                    imap = imaplib.IMAP4_SSL(IMAP_HOST)
                    imap.login(IMAP_USER, IMAP_PASS)
                    imap.select(IMAP_FOLDER)
                    since = (datetime.datetime.utcnow() - datetime.timedelta(days=since_days)).strftime("%d-%b-%Y")
                    status, data = imap.search(None, f'(SINCE {since})')
                    if status != "OK":
                        imap.logout(); return out
                    uids = data[0].split()
                    uids = uids[-limit:]
                    for uid in reversed(uids):
                        status, msg_data = imap.fetch(uid, "(RFC822)")
                        if status != "OK":
                            continue
                        msg = email.message_from_bytes(msg_data[0][1])
                        sender = _decode_header(msg.get("From"))
                        subj   = _decode_header(msg.get("Subject"))
                        body = ""
                        if msg.is_multipart():
                            for part in msg.walk():
                                ctype = part.get_content_type()
                                disp  = part.get("Content-Disposition", "")
                                if ctype == "text/plain" and "attachment" not in (disp or "").lower():
                                    body = part.get_payload(decode=True).decode(errors="ignore"); break
                            if not body:
                                for part in msg.walk():
                                    if part.get_content_type() == "text/html":
                                        html = part.get_payload(decode=True).decode(errors="ignore")
                                        body = re.sub("<[^>]+>", " ", html); break
                        else:
                            if msg.get_content_type() == "text/plain":
                                body = msg.get_payload(decode=True).decode(errors="ignore")
                            elif msg.get_content_type() == "text/html":
                                html = msg.get_payload(decode=True).decode(errors="ignore")
                                body = re.sub("<[^>]+>", " ", html)
                        out.append({
                            "uid": uid.decode() if isinstance(uid, bytes) else str(uid),
                            "from": sender, "subject": subj, "body": (body or "").strip()
                        })
                    imap.logout()
                except Exception as e:
                    st.warning(f"IMAP error: {e}")
                return out

            if st.sidebar.button("Refresh Inbox"):
                st.session_state.pop("inbox_cache", None)
            inbox = st.session_state.get("inbox_cache")
            if inbox is None:
                inbox = fetch_recent_emails(limit=50, since_days=14)
                st.session_state["inbox_cache"] = inbox

            if not inbox:
                st.sidebar.info("No recent emails or IMAP not configured.")
            else:
                labels = [f"{i+1}. {(_decode_header(x['subject']) or '(no subject)')[:80]} — {(_decode_header(x['from']) or '(no from)')[:60]}" for i, x in enumerate(inbox)]
                sel = st.sidebar.selectbox("Select email", range(len(labels)), format_func=lambda i: labels[i])
                chosen = inbox[sel]
                st.sidebar.caption("Selected email:")
                st.sidebar.write(f"**From:** {chosen['from']}")
                st.sidebar.write(f"**Subject:** {chosen['subject']}")
                st.sidebar.write((chosen["body"] or "")[:600] + ("..." if len(chosen["body"] or "")>600 else ""))

                # Draft generator
                def generate_policy_grounded_reply(question_text: str, kb_vs):
                    ctx_docs = []
                    if kb_vs:
                        ctx_docs = kb_vs.as_retriever(search_kwargs={"k": 6}).get_relevant_documents(question_text)
                        context = "\n\n".join(d.page_content for d in ctx_docs)
                    else:
                        context = "(no index loaded)"
                    try:
                        lang = lang_detect((question_text or "")[:4000])
                    except Exception:
                        lang = "en"
                    prompt = f"""Draft a professional, policy-compliant email reply on behalf of {COUNCIL.title()}.
Base ONLY on the provided policy context. Include clear next steps and a friendly closing.
If policy info is missing, say so and suggest the correct contact or form.
Write in the same language as the incoming message (detected: {lang}).

[CONTEXT]
{context}

[INCOMING MESSAGE]
{question_text}

[REPLY DRAFT]
"""
                    draft = get_llm().predict(prompt)
                    # sources
                    srcs = []
                    if ctx_docs:
                        for d in ctx_docs:
                            first = (d.page_content.splitlines() or [""])[0].strip()
                            if first.startswith("[SOURCE]"):
                                s = first.replace("[SOURCE]", "").strip()
                                if s not in srcs:
                                    srcs.append(s)
                    return draft, srcs

                if st.sidebar.button("Generate Draft Reply"):
                    st.session_state["pro_draft"], st.session_state["pro_sources"] = \
                        generate_policy_grounded_reply(chosen["body"] or chosen["subject"], None if 'kb' not in globals() else kb)

                draft = st.session_state.get("pro_draft")
                if draft:
                    st.sidebar.markdown("**Draft Reply**")
                    edited = st.sidebar.text_area("Edit before sending", value=draft, height=220)
                    st.session_state["pro_draft"] = edited

                    srcs = st.session_state.get("pro_sources", [])
                    if srcs:
                        st.sidebar.caption("Sources:")
                        for s in srcs[:6]: st.sidebar.code(s)

                    # prep send
                    to_display = chosen["from"]
                    to_addr = to_display.split("<")[-1].split(">")[0].strip() if "<" in to_display else to_display
                    subj_pref = ("Re: " if not (chosen["subject"] or "").lower().startswith("re:") else "") + (chosen["subject"] or "")
                    subj = st.sidebar.text_input("Subject", value=subj_pref)

                    if st.sidebar.button("Approve & Send"):
                        try:
                            if not (REPLY_FROM and YAGMAIL_PASS):
                                raise RuntimeError("Missing REPLY_FROM or YAGMAIL_PASS")
                            yag = yagmail.SMTP(REPLY_FROM, YAGMAIL_PASS)
                            yag.send(to=to_addr, subject=subj, contents=edited)
                            st.sidebar.success("Email sent.")
                            log_query_to_gcs("pro_email", f"EMAIL:{chosen['subject']}", edited, srcs)
                        except Exception as e:
                            st.sidebar.error(f"Send failed: {e}")

        # Enterprise mini-analytics (only visible to admin)
        if ENTERPRISE_MODE:
            st.sidebar.markdown("### Analytics (Enterprise)")
            with st.sidebar.expander("Last 7 days summary"):
                try:
                    rows = 0
                    for i in range(7):
                        d = (datetime.datetime.utcnow() - datetime.timedelta(days=i)).strftime("%Y%m%d")
                        key = f"logs/{COUNCIL}/queries-{d}.csv"
                        blob = gcs_client().bucket(GCS_BUCKET).blob(key)
                        if blob.exists():
                            txt = blob.download_as_text()
                            rows += max(0, len(txt.splitlines()) - 1)
                    st.write(f"Total logged interactions: **{rows}** (7 days)")
                except Exception as e:
                    st.info(f"Analytics unavailable: {e}")
else:
    st.sidebar.info("Public mode — admin hidden.")

# =========================
#  MAIN MODES
# =========================
def show_extracts_with_links(docs):
    with st.expander("📄 Relevant Extracts (with sources)"):
        for i, d in enumerate(docs, 1):
            src = "Unknown"
            first = (d.page_content.splitlines() or [""])[0].strip()
            if first.startswith("[SOURCE]"):
                src = first.replace("[SOURCE]", "").strip()
            st.markdown(f"**Extract {i}** — _{src}_")
            st.write(d.page_content)
            if src.startswith("gs://"):
                ttl = SIGNED_URL_TTL_MINUTES if ENTERPRISE_MODE else 120
                signed = signed_url_for_blob_path(src, minutes=ttl)
                if signed:
                    st.link_button("View Policy", signed, use_container_width=False, key=f"lnk_{i}_{hash(src)%10**6}")
            st.markdown("---")

kb = None
if PUBLIC_MODE:
    mode = "Knowledge Base"
else:
    mode = st.radio("Choose a data source", ["Knowledge Base", "Upload a File"], horizontal=True)

if mode == "Knowledge Base":
    kb = load_vectorstore()
    if not kb:
        st.info(
            "No FAISS index available yet.\n"
            f"Docs location:  gs://{GCS_BUCKET}/{DOC_PREFIX}\n"
            f"Index location: gs://{GCS_BUCKET}/{INDEX_PREFIX}\n"
            "Tip: set INDEX_ON_START=true or trigger a rebuild from Admin."
        )

q = st.text_input("💬 Ask a question:")

if q:
    if mode == "Knowledge Base":
        if not kb:
            st.error("Knowledge base not loaded.")
        else:
            with st.spinner("🔎 Searching documents…"):
                docs = kb.as_retriever(search_kwargs={"k": 6}).get_relevant_documents(q)
                context = "\n\n".join(d.page_content for d in docs)
                prompt = f"""Answer on behalf of {COUNCIL.title()} Council using ONLY the context.
If the answer is unclear or not covered, say so and suggest the correct next step (link/phone/form).
Keep it concise (4–8 sentences).

[CONTEXT]
{context}

[QUESTION]
{q}
"""
                ans = get_llm().predict(prompt)
            st.subheader("📌 AI Response")
            st.write(ans)
            show_extracts_with_links(docs)
            # sources for logs
            srcs = []
            for d in docs:
                first = (d.page_content.splitlines() or [""])[0].strip()
                if first.startswith("[SOURCE]"):
                    s = first.replace("[SOURCE]", "").strip()
                    if s not in srcs: srcs.append(s)
            log_query_to_gcs("kb", q, ans, srcs)

    else:
        if PUBLIC_MODE:
            st.error("Uploads are disabled in Public mode.")
        else:
            max_bytes = MAX_UPLOAD_MB * 1024 * 1024
            uploaded = st.file_uploader("📄 Upload a document", type=[e.lstrip(".") for e in SUPPORTED_EXTS])
            if not uploaded:
                st.stop()
            if uploaded.size > max_bytes:
                st.error(f"File too large. Limit is {MAX_UPLOAD_MB} MB.")
                st.stop()
            name, data = uploaded.name, uploaded.read()
            body = extract_any(name, data)
            if not body.strip():
                st.error("Couldn’t extract text from that file.")
            else:
                splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
                chunks = splitter.split_text(f"[SOURCE] (uploaded) {name}\n" + body)
                embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
                vs = FAISS.from_texts(chunks, embeddings)
                with st.spinner("📑 Analyzing your document…"):
                    docs = vs.as_retriever(search_kwargs={"k": 6}).get_relevant_documents(q)
                    context = "\n\n".join(d.page_content for d in docs)
                    prompt = f"""Use ONLY the uploaded document context to answer succinctly.
If unsure, say so and suggest a next step. Include a relevant link if present.

[CONTEXT]
{context}

[QUESTION]
{q}
"""
                    ans = get_llm().predict(prompt)
                st.subheader("📌 AI Response")
                st.write(ans)
                show_extracts_with_links(docs)
                log_query_to_gcs("upload", q, ans, [])
                
# Footer
st.markdown(
    f"""
    <hr/>
    <div style='color:#7a7a7a;font-size:12px'>
      Council: <code>{COUNCIL}</code> • Bucket: <code>{GCS_BUCKET}</code> • Docs: <code>{DOC_PREFIX}</code> • Index: <code>{INDEX_PREFIX}</code>
      {"• Data region: <code>"+DATA_REGION+"</code>" if ENTERPRISE_MODE and DATA_REGION else ""}
      <br/>© {datetime.datetime.utcnow().year} IncidentResponder AI
    </div>
    """,
    unsafe_allow_html=True
)
